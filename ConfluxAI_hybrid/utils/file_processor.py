"""
File processor utilities for ConfluxAI Multi-Modal Search Agent
Handles PDF, image, and text file processing
"""
import os
import uuid
import mimetypes
from pathlib import Path
from typing import List, Dict, Any, Optional
import logging
from datetime import datetime

# PDF processing
try:
    import PyPDF2
    from pypdf import PdfReader
    import pdfplumber
    # import tabula  # Commented out - will be imported only when needed
    # import camelot  # Commented out - will be imported only when needed
except ImportError:
    PyPDF2 = None
    PdfReader = None
    pdfplumber = None

# Image processing
try:
    from PIL import Image, ImageEnhance
    import pytesseract
except ImportError:
    Image = None
    pytesseract = None

# Document processing
try:
    from docx import Document
    import openpyxl
    import pandas as pd
    from pptx import Presentation
except ImportError:
    Document = None
    openpyxl = None
    pd = None
    Presentation = None

# Additional file processing
try:
    from bs4 import BeautifulSoup
    import markdown as md
except ImportError:
    BeautifulSoup = None
    md = None

# Computer vision and ML
try:
    import cv2
    import torch
    import torchvision.transforms as transforms
    from transformers.pipelines import pipeline
except ImportError:
    cv2 = None
    torch = None
    transforms = None
    pipeline = None

# Text processing
import re
import json
import numpy as np
from io import BytesIO

from models.schemas import ProcessingResult, ChunkData, ImageAnalysis, PDFAnalysis, TableData
from config.settings import Settings

logger = logging.getLogger(__name__)

class FileProcessor:
    """Handles processing of various file types"""
    
    def __init__(self):
        self.settings = Settings()
        self.supported_types = {
            'application/pdf': self._process_pdf,
            'text/plain': self._process_text,
            'application/vnd.openxmlformats-officedocument.wordprocessingml.document': self._process_docx,
            'application/msword': self._process_doc,
            'image/jpeg': self._process_image,
            'image/png': self._process_image,
            'image/gif': self._process_image,
            'image/bmp': self._process_image,
            'image/tiff': self._process_image,
            'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet': self._process_excel,
            'application/vnd.ms-excel': self._process_excel,
            'text/csv': self._process_csv,
            # Phase 2 additions
            'application/vnd.openxmlformats-officedocument.presentationml.presentation': self._process_powerpoint,
            'application/vnd.ms-powerpoint': self._process_powerpoint,
            'text/html': self._process_html,
            'text/markdown': self._process_markdown,
            'application/json': self._process_json,
            'application/xml': self._process_xml,
            'text/rtf': self._process_rtf,
            # Code files
            'text/x-python': self._process_code,
            'application/javascript': self._process_code,
            'text/x-java-source': self._process_code,
            'text/x-c': self._process_code,
            'text/x-c++': self._process_code,
        }
    
    async def process_file(self, file_path: str, filename: str) -> ProcessingResult:
        """
        Process a file and extract content
        
        Args:
            file_path: Path to the file
            filename: Original filename
            
        Returns:
            ProcessingResult: Processing results
        """
        start_time = datetime.now()
        file_id = str(uuid.uuid4())
        
        try:
            # Detect file type
            content_type, _ = mimetypes.guess_type(filename)
            if not content_type:
                content_type = self._detect_content_type(file_path)
            
            logger.info(f"Processing file {filename} with content type {content_type}")
            
            # Process based on file type
            pdf_analysis = None
            
            if content_type in self.supported_types:
                processor = self.supported_types[content_type]
                
                # Use enhanced PDF processing if available
                if content_type == 'application/pdf' and self.settings.ENABLE_ADVANCED_PDF:
                    text_content, pdf_analysis = await self._process_pdf_advanced(file_path, filename)
                    image_analysis = None
                # Use enhanced image processing if enabled
                elif content_type.startswith('image/') and self.settings.ENABLE_OBJECT_DETECTION:
                    text_content, image_analysis = await self._process_image_advanced(file_path, filename)
                else:
                    text_content, image_analysis = await processor(file_path, filename)
            else:
                # Try to process as text
                text_content, image_analysis = await self._process_text(file_path, filename)
            
            # Create chunks
            chunks = self._create_chunks(text_content, file_id)
            
            processing_time = (datetime.now() - start_time).total_seconds()
            
            return ProcessingResult(
                file_id=file_id,
                content_type=content_type,
                text_content=text_content,
                image_analysis=image_analysis,
                pdf_analysis=pdf_analysis,
                chunks=chunks,
                metadata={
                    'filename': filename,
                    'file_size': os.path.getsize(file_path),
                    'content_type': content_type,
                    'chunks_count': len(chunks)
                },
                processing_time=processing_time
            )
            
        except Exception as e:
            logger.error(f"Error processing file {filename}: {str(e)}")
            raise
    
    def _detect_content_type(self, file_path: str) -> str:
        """Detect content type from file"""
        try:
            with open(file_path, 'rb') as f:
                header = f.read(1024)
            
            # PDF signature
            if header.startswith(b'%PDF'):
                return 'application/pdf'
            
            # Image signatures
            if header.startswith(b'\xff\xd8\xff'):
                return 'image/jpeg'
            elif header.startswith(b'\x89PNG'):
                return 'image/png'
            elif header.startswith(b'GIF8'):
                return 'image/gif'
            elif header.startswith(b'BM'):
                return 'image/bmp'
            
            # Try to decode as text
            try:
                with open(file_path, 'r', encoding='utf-8') as f:
                    f.read(100)
                return 'text/plain'
            except UnicodeDecodeError:
                pass
            
            return 'application/octet-stream'
            
        except Exception:
            return 'application/octet-stream'
    
    async def _process_pdf(self, file_path: str, filename: str) -> tuple[str, Optional[ImageAnalysis]]:
        """Process PDF file"""
        try:
            text_content = ""
            
            if PdfReader:
                # Use pypdf (newer)
                with open(file_path, 'rb') as file:
                    pdf_reader = PdfReader(file)
                    for page in pdf_reader.pages:
                        text_content += page.extract_text() + "\n"
            elif PyPDF2:
                # Fallback to PyPDF2
                with open(file_path, 'rb') as file:
                    pdf_reader = PyPDF2.PdfReader(file)
                    for page in pdf_reader.pages:
                        text_content += page.extract_text() + "\n"
            else:
                raise Exception("No PDF processing library available")
            
            # Clean up text
            text_content = self._clean_text(text_content)
            
            return text_content, None
            
        except Exception as e:
            logger.error(f"Error processing PDF {filename}: {str(e)}")
            raise
    
    async def _process_text(self, file_path: str, filename: str) -> tuple[str, Optional[ImageAnalysis]]:
        """Process text file"""
        try:
            encodings = ['utf-8', 'utf-16', 'latin-1', 'cp1252']
            
            for encoding in encodings:
                try:
                    with open(file_path, 'r', encoding=encoding) as file:
                        text_content = file.read()
                    break
                except UnicodeDecodeError:
                    continue
            else:
                raise Exception("Could not decode text file")
            
            text_content = self._clean_text(text_content)
            return text_content, None
            
        except Exception as e:
            logger.error(f"Error processing text file {filename}: {str(e)}")
            raise
    
    async def _process_docx(self, file_path: str, filename: str) -> tuple[str, Optional[ImageAnalysis]]:
        """Process DOCX file"""
        try:
            if not Document:
                raise Exception("python-docx not available")
            
            doc = Document(file_path)
            text_content = ""
            
            for paragraph in doc.paragraphs:
                text_content += paragraph.text + "\n"
            
            # Extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        text_content += cell.text + " "
                    text_content += "\n"
            
            text_content = self._clean_text(text_content)
            return text_content, None
            
        except Exception as e:
            logger.error(f"Error processing DOCX {filename}: {str(e)}")
            raise
    
    async def _process_doc(self, file_path: str, filename: str) -> tuple[str, Optional[ImageAnalysis]]:
        """Process DOC file (legacy format)"""
        # For DOC files, we might need additional libraries like python-docx2txt
        # For now, treat as text and try to extract what we can
        try:
            return await self._process_text(file_path, filename)
        except Exception as e:
            logger.error(f"Error processing DOC {filename}: {str(e)}")
            raise
    
    async def _process_image(self, file_path: str, filename: str) -> tuple[str, Optional[ImageAnalysis]]:
        """Process image file"""
        try:
            if not Image:
                raise Exception("PIL not available")
            
            # Open and process image
            with Image.open(file_path) as img:
                # Convert to RGB if necessary
                if img.mode != 'RGB':
                    img = img.convert('RGB')
                
                # Resize if too large
                if img.size[0] > self.settings.MAX_IMAGE_SIZE[0] or img.size[1] > self.settings.MAX_IMAGE_SIZE[1]:
                    img.thumbnail(self.settings.MAX_IMAGE_SIZE, Image.Resampling.LANCZOS)
                
                # Extract text using OCR
                text_content = ""
                if pytesseract:
                    try:
                        text_content = pytesseract.image_to_string(img)
                        text_content = self._clean_text(text_content)
                    except Exception as e:
                        logger.warning(f"OCR failed for {filename}: {str(e)}")
                
                # Basic image analysis
                image_analysis = ImageAnalysis(
                    description=f"Image file: {filename}",
                    objects=[],  # Would need more sophisticated ML models
                    text_content=text_content if text_content.strip() else None,
                    features={
                        'size': img.size,
                        'mode': img.mode,
                        'format': img.format
                    },
                    confidence_scores={},
                    bounding_boxes=[],
                    ocr_confidence=0.0 if not text_content.strip() else 0.8
                )
                
                return text_content, image_analysis
                
        except Exception as e:
            logger.error(f"Error processing image {filename}: {str(e)}")
            raise
    
    async def _process_excel(self, file_path: str, filename: str) -> tuple[str, Optional[ImageAnalysis]]:
        """Process Excel file"""
        try:
            if not pd:
                raise Exception("pandas not available")
            
            # Read all sheets
            excel_file = pd.ExcelFile(file_path)
            text_content = ""
            
            for sheet_name in excel_file.sheet_names:
                df = pd.read_excel(file_path, sheet_name=sheet_name)
                text_content += f"Sheet: {sheet_name}\n"
                text_content += df.to_string(index=False) + "\n\n"
            
            text_content = self._clean_text(text_content)
            return text_content, None
            
        except Exception as e:
            logger.error(f"Error processing Excel {filename}: {str(e)}")
            raise
    
    async def _process_csv(self, file_path: str, filename: str) -> tuple[str, Optional[ImageAnalysis]]:
        """Process CSV file"""
        try:
            if not pd:
                raise Exception("pandas not available")
            
            # Try different encodings and separators
            for encoding in ['utf-8', 'latin-1', 'cp1252']:
                for sep in [',', ';', '\t']:
                    try:
                        df = pd.read_csv(file_path, encoding=encoding, sep=sep)
                        text_content = df.to_string(index=False)
                        text_content = self._clean_text(text_content)
                        return text_content, None
                    except Exception:
                        continue
            
            # Fallback to text processing
            return await self._process_text(file_path, filename)
            
        except Exception as e:
            logger.error(f"Error processing CSV {filename}: {str(e)}")
            raise
    
    # Phase 2 Enhancement Methods
    
    async def _process_pdf_advanced(self, file_path: str, filename: str) -> tuple[str, Optional[PDFAnalysis]]:
        """Advanced PDF processing with table and image extraction"""
        try:
            text_content = ""
            tables = []
            images = []
            metadata = {}
            page_count = 0
            
            # First try with pdfplumber for better text and table extraction
            if pdfplumber:
                with pdfplumber.open(file_path) as pdf:
                    page_count = len(pdf.pages)
                    metadata = pdf.metadata or {}
                    
                    for page_num, page in enumerate(pdf.pages):
                        # Extract text with layout preservation
                        page_text = page.extract_text()
                        if page_text:
                            text_content += f"Page {page_num + 1}:\n{page_text}\n\n"
                        
                        # Extract tables if enabled
                        if self.settings.PDF_TABLE_EXTRACTION:
                            page_tables = page.extract_tables()
                            for table_idx, table in enumerate(page_tables):
                                if table and len(table) > 0:
                                    headers = table[0] if table[0] else []
                                    rows = table[1:] if len(table) > 1 else []
                                    
                                    table_data = TableData(
                                        table_id=f"{filename}_p{page_num + 1}_t{table_idx}",
                                        headers=headers,
                                        rows=rows,
                                        confidence=0.8,  # pdfplumber is generally reliable
                                        page_number=page_num + 1,
                                        position=None
                                    )
                                    tables.append(table_data)
                                    
                                    # Add table content to text
                                    text_content += f"\nTable {table_idx + 1} (Page {page_num + 1}):\n"
                                    text_content += f"Headers: {', '.join(headers)}\n"
                                    for row in rows[:5]:  # Limit to first 5 rows
                                        text_content += f"Row: {', '.join(row)}\n"
                                    text_content += "\n"
            
            # Fallback to basic PDF processing
            if not text_content:
                text_content, _ = await self._process_pdf(file_path, filename)
                page_count = 1  # Estimate
            
            # Clean up text
            text_content = self._clean_text(text_content)
            
            pdf_analysis = PDFAnalysis(
                text_content=text_content,
                tables=tables,
                images=images,
                metadata=metadata,
                page_count=page_count,
                layout_preserved=bool(pdfplumber)
            )
            
            return text_content, pdf_analysis
            
        except Exception as e:
            logger.error(f"Advanced PDF processing failed for {filename}: {str(e)}")
            # Fallback to basic processing
            basic_text, _ = await self._process_pdf(file_path, filename)
            return basic_text, None
    
    async def _process_powerpoint(self, file_path: str, filename: str) -> tuple[str, Optional[ImageAnalysis]]:
        """Process PowerPoint file"""
        try:
            if not Presentation:
                raise Exception("python-pptx not available")
            
            prs = Presentation(file_path)
            text_content = ""
            
            for slide_num, slide in enumerate(prs.slides, 1):
                text_content += f"Slide {slide_num}:\n"
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text_content += f"{shape.text}\n"
                
                # Extract notes if available
                if slide.notes_slide and slide.notes_slide.notes_text_frame:
                    notes_text = slide.notes_slide.notes_text_frame.text
                    if notes_text.strip():
                        text_content += f"Notes: {notes_text}\n"
                
                text_content += "\n"
            
            text_content = self._clean_text(text_content)
            return text_content, None
            
        except Exception as e:
            logger.error(f"Error processing PowerPoint {filename}: {str(e)}")
            raise
    
    async def _process_html(self, file_path: str, filename: str) -> tuple[str, Optional[ImageAnalysis]]:
        """Process HTML file"""
        try:
            if not BeautifulSoup:
                raise Exception("beautifulsoup4 not available")
            
            with open(file_path, 'r', encoding='utf-8') as file:
                html_content = file.read()
            
            soup = BeautifulSoup(html_content, 'html.parser')
            
            # Remove script and style elements
            for script in soup(["script", "style"]):
                script.decompose()
            
            # Extract text
            text_content = soup.get_text()
            
            # Extract headings and structure
            headings = []
            for tag in soup.find_all(['h1', 'h2', 'h3', 'h4', 'h5', 'h6']):
                headings.append(f"{tag.name.upper()}: {tag.get_text().strip()}")
            
            if headings:
                text_content = "STRUCTURE:\n" + "\n".join(headings) + "\n\nCONTENT:\n" + text_content
            
            text_content = self._clean_text(text_content)
            return text_content, None
            
        except Exception as e:
            logger.error(f"Error processing HTML {filename}: {str(e)}")
            raise
    
    async def _process_markdown(self, file_path: str, filename: str) -> tuple[str, Optional[ImageAnalysis]]:
        """Process Markdown file"""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                md_content = file.read()
            
            # Convert markdown to HTML and then extract text
            if md:
                html = md.markdown(md_content)
                if BeautifulSoup:
                    soup = BeautifulSoup(html, 'html.parser')
                    text_content = soup.get_text()
                else:
                    text_content = md_content
            else:
                text_content = md_content
            
            text_content = self._clean_text(text_content)
            return text_content, None
            
        except Exception as e:
            logger.error(f"Error processing Markdown {filename}: {str(e)}")
            raise
    
    async def _process_json(self, file_path: str, filename: str) -> tuple[str, Optional[ImageAnalysis]]:
        """Process JSON file"""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                json_data = json.load(file)
            
            # Convert JSON to readable text format
            text_content = f"JSON File: {filename}\n\n"
            text_content += json.dumps(json_data, indent=2, ensure_ascii=False)
            
            text_content = self._clean_text(text_content)
            return text_content, None
            
        except Exception as e:
            logger.error(f"Error processing JSON {filename}: {str(e)}")
            raise
    
    async def _process_xml(self, file_path: str, filename: str) -> tuple[str, Optional[ImageAnalysis]]:
        """Process XML file"""
        try:
            if not BeautifulSoup:
                # Fallback to text processing
                return await self._process_text(file_path, filename)
            
            with open(file_path, 'r', encoding='utf-8') as file:
                xml_content = file.read()
            
            soup = BeautifulSoup(xml_content, 'xml')
            text_content = soup.get_text()
            
            text_content = self._clean_text(text_content)
            return text_content, None
            
        except Exception as e:
            logger.error(f"Error processing XML {filename}: {str(e)}")
            raise
    
    async def _process_rtf(self, file_path: str, filename: str) -> tuple[str, Optional[ImageAnalysis]]:
        """Process RTF file"""
        try:
            # For now, treat as text and extract what we can
            return await self._process_text(file_path, filename)
        except Exception as e:
            logger.error(f"Error processing RTF {filename}: {str(e)}")
            raise
    
    async def _process_code(self, file_path: str, filename: str) -> tuple[str, Optional[ImageAnalysis]]:
        """Process code files"""
        try:
            encodings = ['utf-8', 'utf-16', 'latin-1', 'cp1252']
            
            for encoding in encodings:
                try:
                    with open(file_path, 'r', encoding=encoding) as file:
                        code_content = file.read()
                    break
                except UnicodeDecodeError:
                    continue
            else:
                raise Exception("Could not decode code file")
            
            # Extract file extension for language detection
            file_ext = Path(filename).suffix.lower()
            language_map = {
                '.py': 'Python',
                '.js': 'JavaScript',
                '.java': 'Java',
                '.cpp': 'C++',
                '.c': 'C',
                '.h': 'C Header',
                '.json': 'JSON',
                '.xml': 'XML',
                '.html': 'HTML',
                '.css': 'CSS'
            }
            
            language = language_map.get(file_ext, 'Unknown')
            
            # Add metadata about the code
            text_content = f"Code File: {filename}\n"
            text_content += f"Language: {language}\n"
            text_content += f"File Extension: {file_ext}\n\n"
            text_content += "CODE CONTENT:\n"
            text_content += code_content
            
            text_content = self._clean_text(text_content)
            return text_content, None
            
        except Exception as e:
            logger.error(f"Error processing code file {filename}: {str(e)}")
            raise
    
    async def _process_image_advanced(self, file_path: str, filename: str) -> tuple[str, Optional[ImageAnalysis]]:
        """Advanced image processing with object detection"""
        try:
            if not Image:
                raise Exception("PIL not available")
            
            # Open and process image
            with Image.open(file_path) as img:
                # Convert to RGB if necessary
                if img.mode != 'RGB':
                    img = img.convert('RGB')
                
                # Resize if too large
                original_size = img.size
                if img.size[0] > self.settings.MAX_IMAGE_SIZE[0] or img.size[1] > self.settings.MAX_IMAGE_SIZE[1]:
                    img.thumbnail(self.settings.MAX_IMAGE_SIZE, Image.Resampling.LANCZOS)
                
                # Extract text using OCR with confidence
                text_content = ""
                ocr_confidence = 0.0
                
                if pytesseract:
                    try:
                        # Get text with confidence scores
                        ocr_data = pytesseract.image_to_data(img, output_type=pytesseract.Output.DICT)
                        
                        # Filter text by confidence threshold
                        confident_text = []
                        confidences = []
                        
                        for i, text in enumerate(ocr_data['text']):
                            conf = int(ocr_data['conf'][i])
                            if conf > self.settings.OCR_CONFIDENCE_THRESHOLD * 100 and text.strip():
                                confident_text.append(text)
                                confidences.append(conf)
                        
                        text_content = ' '.join(confident_text)
                        ocr_confidence = sum(confidences) / len(confidences) / 100 if confidences else 0.0
                        
                        text_content = self._clean_text(text_content)
                        
                    except Exception as e:
                        logger.warning(f"OCR failed for {filename}: {str(e)}")
                
                # Object detection (if enabled and available)
                detected_objects = []
                bounding_boxes = []
                confidence_scores = {}
                
                if self.settings.ENABLE_OBJECT_DETECTION and cv2 and torch:
                    try:
                        # Convert PIL to cv2 format
                        img_cv = cv2.cvtColor(np.array(img), cv2.COLOR_RGB2BGR)
                        
                        # Placeholder for object detection
                        # In a real implementation, you would load a YOLO model here
                        # detected_objects, bounding_boxes, confidence_scores = self._detect_objects(img_cv)
                        
                    except Exception as e:
                        logger.warning(f"Object detection failed for {filename}: {str(e)}")
                
                # Enhanced image analysis
                image_analysis = ImageAnalysis(
                    description=f"Image file: {filename} ({original_size[0]}x{original_size[1]})",
                    objects=detected_objects,
                    text_content=text_content if text_content.strip() else None,
                    features={
                        'original_size': original_size,
                        'processed_size': img.size,
                        'mode': img.mode,
                        'format': img.format,
                        'has_text': bool(text_content.strip())
                    },
                    confidence_scores=confidence_scores,
                    bounding_boxes=bounding_boxes,
                    ocr_confidence=ocr_confidence
                )
                
                return text_content, image_analysis
                
        except Exception as e:
            logger.error(f"Error processing image {filename}: {str(e)}")
            raise

    def _detect_objects(self, img_cv) -> tuple[List[str], List[Dict[str, Any]], Dict[str, float]]:
        """Placeholder for object detection implementation"""
        # This would be implemented with a real object detection model
        # For now, return empty results
        return [], [], {}
    
    def _clean_text(self, text: str) -> str:
        """Clean and normalize text"""
        if not text:
            return ""
        
        # Remove excessive whitespace
        text = re.sub(r'\s+', ' ', text)
        
        # Remove control characters
        text = re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f\x7f-\x84\x86-\x9f]', '', text)
        
        # Strip and normalize
        text = text.strip()
        
        return text
    
    def _create_chunks(self, text: str, file_id: str) -> List[ChunkData]:
        """Create text chunks for indexing"""
        if not text:
            return []
        
        chunks = []
        chunk_size = self.settings.CHUNK_SIZE
        overlap = self.settings.CHUNK_OVERLAP
        
        # Split into sentences first
        sentences = re.split(r'[.!?]+', text)
        
        current_chunk = ""
        chunk_index = 0
        
        for sentence in sentences:
            sentence = sentence.strip()
            if not sentence:
                continue
            
            # Check if adding this sentence would exceed chunk size
            if len(current_chunk) + len(sentence) > chunk_size and current_chunk:
                # Create chunk
                chunk_id = f"{file_id}_{chunk_index}"
                chunks.append(ChunkData(
                    chunk_id=chunk_id,
                    content=current_chunk.strip(),
                    file_id=file_id,
                    chunk_index=chunk_index,
                    metadata={'length': len(current_chunk)}
                ))
                
                # Start new chunk with overlap
                if overlap > 0 and len(current_chunk) > overlap:
                    current_chunk = current_chunk[-overlap:] + " " + sentence
                else:
                    current_chunk = sentence
                
                chunk_index += 1
            else:
                if current_chunk:
                    current_chunk += " " + sentence
                else:
                    current_chunk = sentence
        
        # Add final chunk if there's content
        if current_chunk.strip():
            chunk_id = f"{file_id}_{chunk_index}"
            chunks.append(ChunkData(
                chunk_id=chunk_id,
                content=current_chunk.strip(),
                file_id=file_id,
                chunk_index=chunk_index,
                metadata={'length': len(current_chunk)}
            ))
        
        return chunks
